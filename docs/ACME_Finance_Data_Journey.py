#!/usr/bin/env python3
"""Generate ACME Finance Data Journey PowerPoint for interview presentation.

Run: python3 docs/ACME_Finance_Data_Journey.py
Output: docs/ACME_Finance_Data_Journey.pptx
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Brand colors ──────────────────────────────────────────────
DARK_BG     = RGBColor(0x0F, 0x17, 0x2A)    # Deep navy
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)    # Bright blue
ACCENT_TEAL = RGBColor(0x06, 0xB6, 0xD4)    # Cyan/teal
ACCENT_GREEN= RGBColor(0x10, 0xB9, 0x81)    # Green
ACCENT_AMBER= RGBColor(0xF5, 0x9E, 0x0B)    # Amber/orange
ACCENT_RED  = RGBColor(0xEF, 0x44, 0x44)    # Red for alerts
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xCB, 0xD5, 0xE1)
MID_GRAY    = RGBColor(0x64, 0x74, 0x8B)
CARD_BG     = RGBColor(0x1E, 0x29, 0x3B)    # Slightly lighter navy


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_arrow(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_right_arrow(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def set_shape_text(shape, text, font_size=11, color=WHITE, bold=False,
                   alignment=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = alignment
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    shape.text_frame.margin_left = Pt(6)
    shape.text_frame.margin_right = Pt(6)
    shape.text_frame.margin_top = Pt(4)
    shape.text_frame.margin_bottom = Pt(4)


def add_multi_text(shape, lines, font_size=10, color=WHITE, bold_first=False):
    """Add multiple lines to a shape's text frame."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(6)
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        if bold_first and i == 0:
            p.font.bold = True
        p.space_after = Pt(2)


def add_chevron_arrow(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


# ════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide, DARK_BG)

# Accent line
add_shape(slide, Inches(1), Inches(2.2), Inches(1.5), Pt(4), ACCENT_BLUE)

add_text_box(slide, Inches(1), Inches(2.5), Inches(10), Inches(1.2),
             "ACME Finance Platform", font_size=44, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(3.7), Inches(10), Inches(0.8),
             "End-to-End Data Journey — From Source Systems to AI-Powered Analytics",
             font_size=22, color=ACCENT_TEAL, bold=False)
add_text_box(slide, Inches(1), Inches(4.8), Inches(10), Inches(0.6),
             "Architecture Deep Dive  |  Level 200 → 400",
             font_size=16, color=MID_GRAY)

# Tech stack pills
stack_items = ["AWS", "Redshift Serverless", "Apache Iceberg", "Glue ETL",
               "dbt", "Bedrock AI", "Step Functions", "Lambda"]
x_start = Inches(1)
for i, item in enumerate(stack_items):
    pill = add_shape(slide, x_start + Inches(i * 1.35), Inches(5.7),
                     Inches(1.25), Inches(0.38), CARD_BG, ACCENT_BLUE, Pt(1))
    set_shape_text(pill, item, font_size=9, color=ACCENT_BLUE, bold=True)

add_text_box(slide, Inches(1), Inches(6.6), Inches(10), Inches(0.4),
             "Uday Nag  |  Data & Analytics Platform Engineer",
             font_size=13, color=MID_GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.7),
             "Agenda", font_size=36, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.05), Inches(1), Pt(3), ACCENT_BLUE)

agenda_items = [
    ("01", "High-Level Data Flow", "Source → Lake → Warehouse → Consumption\nThe 10,000-foot view", ACCENT_BLUE),
    ("02", "Source Systems", "ERP, CRM, EPM — what data, how generated,\nhow reconciled", ACCENT_TEAL),
    ("03", "Data Lake (S3 + Iceberg)", "Medallion architecture: Raw → Curated\nGlue ETL with ACID guarantees", ACCENT_GREEN),
    ("04", "Transformation (dbt)", "3-layer modeling: Staging → Intermediate → Marts\n23 models, 41 tests, full lineage", ACCENT_AMBER),
    ("05", "Orchestration", "Step Functions pipeline: Crawlers → ETL →\ndbt → Consumption-ready", RGBColor(0xA7, 0x8B, 0xFA)),
    ("06", "Consumption Layer", "React dashboards, Bedrock AI Agent,\ndriver-based ARR forecasting", ACCENT_RED),
]

for i, (num, title, desc, color) in enumerate(agenda_items):
    y = Inches(1.5) + Inches(i * 0.9)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), y, Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    circle.shadow.inherit = False
    set_shape_text(circle, num, font_size=16, color=WHITE, bold=True)
    # Title + description
    add_text_box(slide, Inches(1.8), y - Inches(0.05), Inches(5), Inches(0.35),
                 title, font_size=18, color=WHITE, bold=True)
    add_text_box(slide, Inches(1.8), y + Inches(0.30), Inches(6), Inches(0.55),
                 desc, font_size=11, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 3 — HIGH-LEVEL DATA FLOW (Level 200)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(8), Inches(0.6),
             "High-Level Data Flow", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.5), Inches(0.9), Inches(1), Pt(3), ACCENT_BLUE)
add_text_box(slide, Inches(0.5), Inches(1.0), Inches(6), Inches(0.4),
             "LEVEL 200  |  The 10,000-foot view", font_size=12, color=ACCENT_BLUE, bold=True)

# ── Layer 1: Source Systems ──
y_base = Inches(1.7)
layer_label = add_text_box(slide, Inches(0.3), y_base, Inches(1.2), Inches(0.3),
                            "SOURCE", font_size=9, color=ACCENT_BLUE, bold=True)

sources = [("ERP", "GL, AR/AP\nCustomers\nFixed Assets"), ("CRM", "Accounts, Opps\nARR Movement\nPipeline"), ("EPM", "Budgets\nForecasts\nHeadcount")]
for i, (name, desc) in enumerate(sources):
    box = add_shape(slide, Inches(1.6 + i*2.5), y_base, Inches(2.2), Inches(0.95), CARD_BG, ACCENT_BLUE)
    add_multi_text(box, [name, desc], font_size=10, color=WHITE, bold_first=True)

# Arrow down
for i in range(3):
    add_arrow(slide, Inches(2.5 + i*2.5), y_base + Inches(1.05), Inches(0.3), Inches(0.35), ACCENT_BLUE)

# ── Layer 2: Data Lake ──
y2 = y_base + Inches(1.5)
add_text_box(slide, Inches(0.3), y2, Inches(1.2), Inches(0.3),
             "DATA LAKE", font_size=9, color=ACCENT_TEAL, bold=True)

raw_box = add_shape(slide, Inches(1.6), y2, Inches(3.5), Inches(0.85), CARD_BG, ACCENT_TEAL)
add_multi_text(raw_box, ["RAW ZONE (Bronze)", "S3 Parquet  |  Hive-style partitions", "erp/ crm/ epm/ prefixes"], font_size=10, bold_first=True)

arrow_r = add_right_arrow(slide, Inches(5.3), y2 + Inches(0.2), Inches(0.6), Inches(0.4), ACCENT_TEAL)

curated_box = add_shape(slide, Inches(6.1), y2, Inches(3.5), Inches(0.85), CARD_BG, ACCENT_GREEN)
add_multi_text(curated_box, ["CURATED ZONE (Silver)", "Apache Iceberg v2  |  ACID writes", "Glue ETL (PySpark createOrReplace)"], font_size=10, bold_first=True)

# Arrow down
add_arrow(slide, Inches(7.6), y2 + Inches(0.95), Inches(0.3), Inches(0.35), ACCENT_GREEN)

# ── Layer 3: Warehouse ──
y3 = y2 + Inches(1.45)
add_text_box(slide, Inches(0.3), y3, Inches(1.4), Inches(0.3),
             "WAREHOUSE", font_size=9, color=ACCENT_AMBER, bold=True)

wh_box = add_shape(slide, Inches(1.6), y3, Inches(8), Inches(0.85), CARD_BG, ACCENT_AMBER)
add_multi_text(wh_box, [
    "Redshift Serverless  +  dbt (23 models, 41 tests)",
    "Staging (schema conform)  →  Intermediate (business logic)  →  Marts (consumption-ready)",
    "mart_pl  |  fct_arr  |  fct_revenue  |  fct_expense  |  mart_ar_aging  |  6 dim tables"
], font_size=10, bold_first=True)

# Arrow down
add_arrow(slide, Inches(5.3), y3 + Inches(0.95), Inches(0.3), Inches(0.35), ACCENT_AMBER)

# ── Layer 4: Consumption ──
y4 = y3 + Inches(1.45)
add_text_box(slide, Inches(0.3), y4, Inches(1.4), Inches(0.3),
             "CONSUME", font_size=9, color=ACCENT_RED, bold=True)

apps = [
    ("React Dashboard", "P&L, ARR, Forecast\nAnomaly Detection\nCloudFront CDN"),
    ("Bedrock AI Agent", "NL → SQL → Answer\nVariance RCA\n6 Action Groups"),
    ("Forecast Lambda", "Driver-based ARR\nTier cohort model\nScenario what-if"),
]
for i, (name, desc) in enumerate(apps):
    color = [ACCENT_RED, RGBColor(0xA7, 0x8B, 0xFA), ACCENT_AMBER][i]
    box = add_shape(slide, Inches(1.6 + i*2.8), y4, Inches(2.5), Inches(0.85), CARD_BG, color)
    add_multi_text(box, [name, desc], font_size=10, color=WHITE, bold_first=True)

# Side annotation
add_text_box(slide, Inches(10.2), y_base, Inches(2.8), Inches(5.5),
    "KEY POINTS\n\n"
    "• Medallion architecture\n  (Bronze → Silver → Gold)\n\n"
    "• Zero-copy via Redshift\n  Spectrum (external schema)\n\n"
    "• Glue Catalog as unified\n  metadata store\n\n"
    "• Step Functions orchestrates\n  end-to-end pipeline\n\n"
    "• IAM auth everywhere\n  (no stored passwords)",
    font_size=10, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 4 — SOURCE SYSTEMS DEEP DIVE
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(8), Inches(0.6),
             "Source Systems", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.5), Inches(0.9), Inches(1), Pt(3), ACCENT_BLUE)
add_text_box(slide, Inches(0.5), Inches(1.0), Inches(6), Inches(0.4),
             "LEVEL 300  |  What data, how structured, how reconciled", font_size=12, color=ACCENT_BLUE, bold=True)

# Three source system cards
systems = [
    ("ERP System", "(NetSuite / SAP-like)", ACCENT_BLUE,
     ["13 tables, ~800K rows",
      "GL Journals (header + line)",
      "AR Invoices & Receipts",
      "AP Invoices & Payments",
      "Customers, Vendors",
      "Fixed Assets + Depreciation",
      "Chart of Accounts",
      "Entities (US, EMEA, APAC)",
      "",
      "Format: Parquet (non-partitioned)",
      "Also: CSV for Postgres COPY"]),
    ("CRM System", "(Salesforce)", ACCENT_TEAL,
     ["6 tables, ~10K rows",
      "Accounts (= ERP customers)",
      "Contacts",
      "Opportunities + Line Items",
      "ARR Movements (partitioned)",
      "Pipeline Snapshots (partitioned)",
      "",
      "Key: account_id = customer_id",
      "(shared UUIDs for reconciliation)",
      "",
      "Hive partitions: period_yyyymm"]),
    ("EPM System", "(Anaplan / Adaptive)", ACCENT_GREEN,
     ["5 tables, ~6K rows",
      "Budget Versions",
      "Forecast Versions",
      "Plan Lines (partitioned)",
      "Headcount Plan",
      "Driver Assumptions",
      "",
      "Account-level planning:",
      "budget vs. forecast vs. actual",
      "",
      "Hive partitions: period_yyyymm"]),
]

for i, (name, sub, color, items) in enumerate(systems):
    x = Inches(0.5 + i * 4.2)
    # Header
    header = add_shape(slide, x, Inches(1.6), Inches(3.9), Inches(0.65), color)
    set_shape_text(header, f"{name}\n{sub}", font_size=14, color=WHITE, bold=True)
    # Body
    body = add_shape(slide, x, Inches(2.25), Inches(3.9), Inches(3.4), CARD_BG, color)
    add_multi_text(body, items, font_size=10, color=LIGHT_GRAY)

# Bottom callout — reconciliation
callout = add_shape(slide, Inches(0.5), Inches(5.9), Inches(12.3), Inches(1.2), CARD_BG, ACCENT_AMBER)
add_multi_text(callout, [
    "DATA RECONCILIATION & INTEGRITY",
    "• CRM account_id = ERP customer_id (shared UUIDs generated in Phase 2B, carried through Phase 2C)",
    "• Every GL journal balances: SUM(debits) = SUM(credits) — enforced at generation time",
    "• Revenue recognition: ratable across service periods, capped at fiscal year boundary (Jan 31 FY-end)",
    "• 3 fiscal years generated: FY23 ($1.6B) → FY24 ($2.0B) → FY25 ($2.3B) — realistic SaaS growth trajectory",
    "• Seeded anomalies: aged AR (>90d), S&M overspend (2x spike), asset disposal losses, stalled pipeline deals"
], font_size=10, color=LIGHT_GRAY, bold_first=True)


# ════════════════════════════════════════════════════════════════
# SLIDE 5 — DATA LAKE: S3 + ICEBERG
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(8), Inches(0.6),
             "Data Lake — S3 + Apache Iceberg", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.5), Inches(0.9), Inches(1), Pt(3), ACCENT_TEAL)
add_text_box(slide, Inches(0.5), Inches(1.0), Inches(8), Inches(0.4),
             "LEVEL 400  |  Raw zone layout, Glue ETL, Iceberg writes", font_size=12, color=ACCENT_TEAL, bold=True)

# Left: Raw Zone
raw_header = add_shape(slide, Inches(0.5), Inches(1.7), Inches(5.5), Inches(0.5), ACCENT_TEAL)
set_shape_text(raw_header, "RAW ZONE (Bronze)  —  s3://acme-lake-dev-raw-...", font_size=12, color=WHITE, bold=True)

raw_body = add_shape(slide, Inches(0.5), Inches(2.2), Inches(5.5), Inches(3.0), CARD_BG, ACCENT_TEAL)
add_multi_text(raw_body, [
    "S3 Layout (preserves source system boundaries):",
    "",
    "  erp/",
    "    gl_journal_header/part-0.parquet",
    "    gl_journal_line/part-0.parquet",
    "    customer/part-0.parquet",
    "    ar_invoice/part-0.parquet  ... (13 tables)",
    "",
    "  crm/",
    "    account/part-0.parquet",
    "    arr_movement/",
    "      period_yyyymm=202302/part-0.parquet",
    "      period_yyyymm=202303/part-0.parquet",
    "",
    "  epm/",
    "    plan_line/",
    "      period_yyyymm=202302/part-0.parquet",
], font_size=9, color=LIGHT_GRAY, bold_first=True)

# Arrow
add_right_arrow(slide, Inches(6.2), Inches(3.2), Inches(0.5), Inches(0.35), ACCENT_GREEN)

# Right: Curated Zone
cur_header = add_shape(slide, Inches(6.9), Inches(1.7), Inches(5.9), Inches(0.5), ACCENT_GREEN)
set_shape_text(cur_header, "CURATED ZONE (Silver)  —  Apache Iceberg v2", font_size=12, color=WHITE, bold=True)

cur_body = add_shape(slide, Inches(6.9), Inches(2.2), Inches(5.9), Inches(3.0), CARD_BG, ACCENT_GREEN)
add_multi_text(cur_body, [
    "Glue ETL Job (PySpark):",
    "",
    "1. Read via DynamicFrame (handles partitions)",
    "2. Add audit columns:",
    "     _ingest_date = current_date()",
    "     _ingest_run_id = unique run UUID",
    "3. Write as Iceberg:",
    "     df.writeTo(target_table)",
    '       .tableProperty("format-version", "2")',
    "       .createOrReplace()",
    "",
    "Why Iceberg v2?",
    "  • ACID guarantees (atomic metadata swap)",
    "  • Schema evolution (add columns safely)",
    "  • Time travel (query historical snapshots)",
    "  • Redshift Spectrum compatible (zero-copy)"
], font_size=9, color=LIGHT_GRAY, bold_first=True)

# Bottom: Glue Crawlers
crawl_header = add_shape(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.4), RGBColor(0xA7, 0x8B, 0xFA))
set_shape_text(crawl_header, "GLUE CRAWLERS  —  4 crawlers populate Glue Data Catalog (unified metadata store)", font_size=11, color=WHITE, bold=True)

crawlers = [
    ("crawler-erp", "raw_erp_dev", "Parquet schemas"),
    ("crawler-crm", "raw_crm_dev", "Hive partitions"),
    ("crawler-epm", "raw_epm_dev", "Hive partitions"),
    ("crawler-curated", "curated_dev", "Iceberg metadata"),
]
for i, (name, db, detail) in enumerate(crawlers):
    c_box = add_shape(slide, Inches(0.5 + i*3.15), Inches(6.0), Inches(2.9), Inches(0.7), CARD_BG, RGBColor(0xA7, 0x8B, 0xFA))
    add_multi_text(c_box, [name, f"→ {db}", detail], font_size=9, color=LIGHT_GRAY, bold_first=True)


# ════════════════════════════════════════════════════════════════
# SLIDE 6 — STEP FUNCTIONS ORCHESTRATION
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(8), Inches(0.6),
             "Pipeline Orchestration — Step Functions", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.5), Inches(0.9), Inches(1), Pt(3), RGBColor(0xA7, 0x8B, 0xFA))
add_text_box(slide, Inches(0.5), Inches(1.0), Inches(8), Inches(0.4),
             "LEVEL 400  |  End-to-end pipeline from trigger to consumption-ready data", font_size=12,
             color=RGBColor(0xA7, 0x8B, 0xFA), bold=True)

# Step 1: EventBridge trigger
step_y = Inches(1.8)
s1 = add_shape(slide, Inches(0.8), step_y, Inches(2.2), Inches(0.7), CARD_BG, ACCENT_BLUE)
add_multi_text(s1, ["EventBridge Rule", "06:00 UTC daily (configurable)"], font_size=10, bold_first=True)

add_right_arrow(slide, Inches(3.15), step_y + Inches(0.15), Inches(0.4), Inches(0.35), ACCENT_BLUE)

# Step 2: Step Functions
s2 = add_shape(slide, Inches(3.7), step_y, Inches(2.5), Inches(0.7), ACCENT_BLUE)
set_shape_text(s2, "Step Functions\nacme-finance-dev-refresh", font_size=11, color=WHITE, bold=True)

# Parallel branches
add_arrow(slide, Inches(4.2), step_y + Inches(0.8), Inches(0.25), Inches(0.35), ACCENT_TEAL)
add_arrow(slide, Inches(4.85), step_y + Inches(0.8), Inches(0.25), Inches(0.35), ACCENT_TEAL)
add_arrow(slide, Inches(5.5), step_y + Inches(0.8), Inches(0.25), Inches(0.35), ACCENT_TEAL)

step_y2 = step_y + Inches(1.3)
parallel_label = add_text_box(slide, Inches(3.0), step_y2 - Inches(0.2), Inches(0.7), Inches(0.3),
                               "Parallel", font_size=8, color=ACCENT_TEAL, bold=True)

branches = [("Glue ETL\nraw_erp → curated\n(13 tables)", 3.3),
            ("Glue ETL\nraw_crm → curated\n(6 tables)", 5.0),
            ("Glue ETL\nraw_epm → curated\n(5 tables)", 6.7)]
for text, x in branches:
    b = add_shape(slide, Inches(x), step_y2, Inches(1.5), Inches(0.9), CARD_BG, ACCENT_TEAL)
    add_multi_text(b, text.split("\n"), font_size=9, bold_first=True)

# Converge
step_y3 = step_y2 + Inches(1.1)
add_arrow(slide, Inches(3.95), step_y3, Inches(0.25), Inches(0.25), ACCENT_GREEN)
add_arrow(slide, Inches(5.65), step_y3, Inches(0.25), Inches(0.25), ACCENT_GREEN)
add_arrow(slide, Inches(7.35), step_y3, Inches(0.25), Inches(0.25), ACCENT_GREEN)

step_y4 = step_y3 + Inches(0.4)
s3 = add_shape(slide, Inches(3.5), step_y4, Inches(3.0), Inches(0.6), CARD_BG, ACCENT_GREEN)
add_multi_text(s3, ["Glue Crawler (curated)", "Detect new Iceberg table versions"], font_size=10, bold_first=True)

add_arrow(slide, Inches(4.85), step_y4 + Inches(0.7), Inches(0.25), Inches(0.3), ACCENT_AMBER)

step_y5 = step_y4 + Inches(1.1)
s4 = add_shape(slide, Inches(3.5), step_y5, Inches(3.0), Inches(0.6), CARD_BG, ACCENT_AMBER)
add_multi_text(s4, ["dbt run", "23 models  |  41 tests  |  staging → marts"], font_size=10, bold_first=True)

add_arrow(slide, Inches(4.85), step_y5 + Inches(0.7), Inches(0.25), Inches(0.3), ACCENT_GREEN)

step_y6 = step_y5 + Inches(1.1)
s5 = add_shape(slide, Inches(3.5), step_y6, Inches(3.0), Inches(0.5), ACCENT_GREEN)
set_shape_text(s5, "DATA READY", font_size=14, color=WHITE, bold=True)

# Right column: annotations
ann_x = Inches(9.0)
annotations = [
    ("Idempotent", "Full-refresh (createOrReplace)\nRe-run any branch safely\nIceberg handles metadata atomically", ACCENT_TEAL),
    ("Observable", "CloudWatch logs per Glue job\nStep Functions visual workflow\nGlue crawler run metrics", ACCENT_GREEN),
    ("Scalable", "Glue auto-scales workers\nRedshift 8–32 RPU range\nServerless — pay only when active", ACCENT_AMBER),
    ("Tested", "dbt: 41 tests on every run\nGL balance assertions\nP&L identity checks", ACCENT_RED),
]
for i, (title, desc, color) in enumerate(annotations):
    ay = Inches(1.8 + i * 1.35)
    ann = add_shape(slide, ann_x, ay, Inches(3.8), Inches(1.15), CARD_BG, color)
    add_multi_text(ann, [title] + desc.split("\n"), font_size=9, color=LIGHT_GRAY, bold_first=True)


# ════════════════════════════════════════════════════════════════
# SLIDE 7 — dbt TRANSFORMATION LAYERS
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(8), Inches(0.6),
             "dbt Transformation — 3-Layer Modeling", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.5), Inches(0.9), Inches(1), Pt(3), ACCENT_AMBER)
add_text_box(slide, Inches(0.5), Inches(1.0), Inches(8), Inches(0.4),
             "LEVEL 400  |  Staging → Intermediate → Marts on Redshift Serverless", font_size=12,
             color=ACCENT_AMBER, bold=True)

# Layer 1: Staging
l1_header = add_shape(slide, Inches(0.5), Inches(1.7), Inches(3.8), Inches(0.45), ACCENT_BLUE)
set_shape_text(l1_header, "STAGING (analytics_dev_staging)", font_size=12, color=WHITE, bold=True)
l1_body = add_shape(slide, Inches(0.5), Inches(2.15), Inches(3.8), Inches(2.4), CARD_BG, ACCENT_BLUE)
add_multi_text(l1_body, [
    "Purpose: 1:1 schema-conform views",
    "",
    "stg_erp__gl_journal_header",
    "stg_erp__gl_journal_line",
    "stg_erp__ar_invoice",
    "stg_erp__customer",
    "stg_crm__account",
    "stg_crm__arr_movement",
    "stg_crm__opportunity",
    "stg_epm__plan_line",
    "",
    "Key: account_id → customer_id rename",
    "(enables cross-system joins)"
], font_size=9, color=LIGHT_GRAY, bold_first=True)

# Arrow
add_right_arrow(slide, Inches(4.5), Inches(3.0), Inches(0.4), Inches(0.3), ACCENT_BLUE)

# Layer 2: Intermediate
l2_header = add_shape(slide, Inches(5.1), Inches(1.7), Inches(3.5), Inches(0.45), ACCENT_TEAL)
set_shape_text(l2_header, "INTERMEDIATE (analytics_dev_int)", font_size=12, color=WHITE, bold=True)
l2_body = add_shape(slide, Inches(5.1), Inches(2.15), Inches(3.5), Inches(2.4), CARD_BG, ACCENT_TEAL)
add_multi_text(l2_body, [
    "Purpose: Business logic joins",
    "",
    "int_gl_entries_enriched",
    "  GL lines + COA + cost center",
    "  + entity mapping",
    "  + account_type, pnl_rollup",
    "",
    "int_revenue_monthly",
    "  SUM(amount) WHERE is_revenue",
    "  GROUP BY entity, segment, period",
    "",
    "int_expense_monthly",
    "  SUM(amount) WHERE is_expense",
    "  GROUP BY entity, rollup, period",
], font_size=9, color=LIGHT_GRAY, bold_first=True)

# Arrow
add_right_arrow(slide, Inches(8.8), Inches(3.0), Inches(0.4), Inches(0.3), ACCENT_TEAL)

# Layer 3: Marts
l3_header = add_shape(slide, Inches(9.4), Inches(1.7), Inches(3.5), Inches(0.45), ACCENT_GREEN)
set_shape_text(l3_header, "MARTS (analytics_dev_marts)", font_size=12, color=WHITE, bold=True)
l3_body = add_shape(slide, Inches(9.4), Inches(2.15), Inches(3.5), Inches(2.4), CARD_BG, ACCENT_GREEN)
add_multi_text(l3_body, [
    "Purpose: Consumption-ready tables",
    "",
    "mart_pl — P&L by entity/period",
    "fct_arr — ARR waterfall (customer)",
    "fct_revenue — by entity/segment",
    "fct_expense — by cost center",
    "fct_gl_entries — atomic GL lines",
    "mart_ar_aging — open invoices",
    "",
    "dim_customer, dim_account",
    "dim_entity, dim_cost_center",
    "dim_date (fiscal calendar spine)",
    "",
    "Total: 23 models"
], font_size=9, color=LIGHT_GRAY, bold_first=True)

# Bottom: Test framework
test_header = add_shape(slide, Inches(0.5), Inches(4.9), Inches(12.3), Inches(0.4), ACCENT_AMBER)
set_shape_text(test_header, "dbt TEST FRAMEWORK  —  41 tests running on every dbt run", font_size=11, color=WHITE, bold=True)

tests = [
    ("Uniqueness", "All primary keys across\ndim_* and fct_*"),
    ("Not-Null", "Primary keys +\ncritical amounts"),
    ("Accepted Values", "account_type: asset,\nliability, equity,\nrevenue, expense"),
    ("GL Balance", "Custom: every journal\nSUM(debits) =\nSUM(credits)"),
    ("P&L Identity", "Custom: operating_income\n= gross_profit\n- total_opex"),
]
for i, (name, desc) in enumerate(tests):
    t_box = add_shape(slide, Inches(0.5 + i * 2.5), Inches(5.4), Inches(2.3), Inches(1.1), CARD_BG, ACCENT_AMBER)
    add_multi_text(t_box, [name, desc], font_size=9, color=LIGHT_GRAY, bold_first=True)

# Connection note
add_text_box(slide, Inches(0.5), Inches(6.7), Inches(12), Inches(0.4),
    "Redshift Spectrum external schema reads Iceberg tables directly from S3 — zero data copying. Spectrum pushes predicates to S3 for efficient scans.",
    font_size=10, color=MID_GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 8 — CONSUMPTION: DASHBOARDS + AI
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(8), Inches(0.6),
             "Consumption Layer — Dashboards & AI", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.5), Inches(0.9), Inches(1), Pt(3), ACCENT_RED)
add_text_box(slide, Inches(0.5), Inches(1.0), Inches(8), Inches(0.4),
             "LEVEL 300  |  How end users interact with the data", font_size=12, color=ACCENT_RED, bold=True)

# Card 1: React Dashboard
c1 = add_shape(slide, Inches(0.5), Inches(1.7), Inches(4.0), Inches(0.5), ACCENT_BLUE)
set_shape_text(c1, "React Dashboard (CloudFront)", font_size=13, color=WHITE, bold=True)
c1b = add_shape(slide, Inches(0.5), Inches(2.2), Inches(4.0), Inches(3.2), CARD_BG, ACCENT_BLUE)
add_multi_text(c1b, [
    "React 18 + TypeScript + Vite + Recharts",
    "",
    "Pages:",
    "  /pl — P&L statement by entity/quarter",
    "  /arr — ARR waterfall chart",
    "  /ar-aging — Aging buckets dashboard",
    "  /forecast — Interactive forecast charts",
    "  /anomalies — Severity-badge anomaly list",
    "  /chat — AI Analyst (streaming SSE)",
    "  /commentary — CFO commentary generator",
    "  /board-pack — PDF board pack export",
    "",
    "Hosting: S3 + CloudFront (OAC + sigv4)",
    "API: Lambda Function URL + LWA → FastAPI",
    "SSE: RESPONSE_STREAM mode for /chat"
], font_size=9, color=LIGHT_GRAY, bold_first=True)

# Card 2: Bedrock Agent
c2 = add_shape(slide, Inches(4.8), Inches(1.7), Inches(4.0), Inches(0.5), RGBColor(0xA7, 0x8B, 0xFA))
set_shape_text(c2, "Bedrock AI Agent (Claude Sonnet)", font_size=13, color=WHITE, bold=True)
c2b = add_shape(slide, Inches(4.8), Inches(2.2), Inches(4.0), Inches(3.2), CARD_BG, RGBColor(0xA7, 0x8B, 0xFA))
add_multi_text(c2b, [
    "6 Action Groups (Lambda functions):",
    "",
    "1. QueryFinanceData — NL → SQL → Answer",
    '   "What was Q3 revenue by entity?"',
    "",
    "2. ForecastMetrics — Driver-based ARR",
    '   "Forecast revenue for 12 months"',
    "",
    "3. VarianceRCA — Actuals vs. Budget",
    '   "Why did we miss Q2 budget?"',
    "",
    "4. WhatIfSimulation — Scenario modeling",
    '   "What if R&D drops 15%?"',
    "",
    "5. AnomalyDetection — Health scanner",
    "6. MetricGlossary — KPI definitions"
], font_size=9, color=LIGHT_GRAY, bold_first=True)

# Card 3: API + Lambda
c3 = add_shape(slide, Inches(9.1), Inches(1.7), Inches(3.7), Inches(0.5), ACCENT_AMBER)
set_shape_text(c3, "API Layer (FastAPI on Lambda)", font_size=13, color=WHITE, bold=True)
c3b = add_shape(slide, Inches(9.1), Inches(2.2), Inches(3.7), Inches(3.2), CARD_BG, ACCENT_AMBER)
add_multi_text(c3b, [
    "Lambda Web Adapter architecture:",
    "",
    "Lambda invocation",
    "  → LWA extension intercepts",
    "  → HTTP to uvicorn :8000",
    "  → FastAPI handles request",
    "",
    "Endpoints:",
    "  /metrics/* — dashboard data",
    "  /chat/stream — SSE to Bedrock",
    "  /commentary — AI-generated notes",
    "  /boardpack — PDF generation",
    "",
    "Auth: IAM-based (Redshift Data API)",
    "Auto-provisioned user:",
    "  IAMR:acme-finance-dev-*-lambda"
], font_size=9, color=LIGHT_GRAY, bold_first=True)

# Bottom: AgentCore Memory
mem = add_shape(slide, Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.9), CARD_BG, RGBColor(0xA7, 0x8B, 0xFA))
add_multi_text(mem, [
    "AgentCore Memory (Semantic) — Cross-session conversation memory",
    "1. User asks question → retrieve top-5 semantically relevant past Q&A pairs",
    "2. Memory context prepended to prompt → Bedrock Agent responds with historical awareness",
    "3. New Q&A stored back (fire-and-forget) → builds institutional knowledge over time"
], font_size=10, color=LIGHT_GRAY, bold_first=True)


# ════════════════════════════════════════════════════════════════
# SLIDE 9 — DRIVER-BASED ARR FORECASTING
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.6),
             "Driver-Based ARR Forecasting", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.5), Inches(0.9), Inches(1), Pt(3), ACCENT_GREEN)
add_text_box(slide, Inches(0.5), Inches(1.0), Inches(8), Inches(0.4),
             "LEVEL 400  |  SaaS cohort model with scenario overrides", font_size=12, color=ACCENT_GREEN, bold=True)

# ARR Bridge Formula
formula_box = add_shape(slide, Inches(0.5), Inches(1.7), Inches(7.5), Inches(2.5), CARD_BG, ACCENT_GREEN)
add_multi_text(formula_box, [
    "ARR BRIDGE FORMULA (per tier, per month)",
    "",
    "Starting_ARR    = prior month's Ending_ARR",
    "Churn           = -(Starting_ARR x monthly_churn_rate)",
    "Contraction     = -(Starting_ARR x monthly_contraction_rate)",
    "Expansion       = Starting_ARR x monthly_expansion_rate",
    "New Logo        = blended monthly ACV",
    "                  (pipeline for months 1-3, run-rate for 4-12)",
    "Ending_ARR      = Starting + Churn + Contraction + Expansion + New",
    "",
    "Monthly Revenue = (Total Ending_ARR / 12) x (1 + nonsub_uplift)",
    "Expenses        = Revenue x trailing OpEx ratios (COGS, S&M, R&D, G&A)",
    "Operating Income = Gross Profit - Total OpEx"
], font_size=10, color=LIGHT_GRAY, bold_first=True)

# Tier metrics table
tier_header = add_shape(slide, Inches(8.3), Inches(1.7), Inches(4.5), Inches(0.45), ACCENT_AMBER)
set_shape_text(tier_header, "TIER DYNAMICS (from trailing 4 quarters)", font_size=11, color=WHITE, bold=True)

tier_labels = ["Metric", "Enterprise", "Commercial", "SMB"]
tier_data = [
    ["NRR", "107%", "89%", "49%"],
    ["GRR", "96%", "79%", "45%"],
    ["Churn Rate", "~0%", "17%", "49%"],
    ["Expansion", "11%", "9%", "4%"],
]

# Table header row
for j, label in enumerate(tier_labels):
    col_x = Inches(8.3 + j * 1.12)
    hdr = add_shape(slide, col_x, Inches(2.15), Inches(1.1), Inches(0.35), ACCENT_AMBER)
    set_shape_text(hdr, label, font_size=9, color=WHITE, bold=True)

# Table data rows
for i, row in enumerate(tier_data):
    for j, val in enumerate(row):
        col_x = Inches(8.3 + j * 1.12)
        cell_color = CARD_BG
        cell = add_shape(slide, col_x, Inches(2.5 + i * 0.35), Inches(1.1), Inches(0.35), cell_color,
                        ACCENT_AMBER if j == 0 else MID_GRAY)
        set_shape_text(cell, val, font_size=9, color=WHITE if j == 0 else LIGHT_GRAY,
                      bold=j == 0)

# Scenario overrides
scenario_header = add_shape(slide, Inches(8.3), Inches(4.0), Inches(4.5), Inches(0.45), ACCENT_RED)
set_shape_text(scenario_header, "SCENARIO OVERRIDES (what-if)", font_size=11, color=WHITE, bold=True)

scenarios = [
    ("churn_pct_multiplier", "0.5 = halve churn\n→ +4.5% revenue"),
    ("expansion_pct_multiplier", "1.5 = 50% more upsell"),
    ("new_logo_pct_change", "+50 = 50% more bookings\n→ +23.5% revenue"),
]
for i, (param, effect) in enumerate(scenarios):
    s_box = add_shape(slide, Inches(8.3), Inches(4.55 + i * 0.8), Inches(4.5), Inches(0.7), CARD_BG, ACCENT_RED)
    add_multi_text(s_box, [param, effect], font_size=9, color=LIGHT_GRAY, bold_first=True)

# Data flow at bottom
flow_header = add_shape(slide, Inches(0.5), Inches(4.5), Inches(7.5), Inches(0.4), ACCENT_BLUE)
set_shape_text(flow_header, "DATA FLOW  —  3 Redshift queries → Roll Forward → P&L Projection", font_size=11, color=WHITE, bold=True)

queries = [
    ("Q1: fct_arr (T4Q)", "Churn, contraction,\nexpansion rates\nby tier"),
    ("Q2: fct_arr + opp", "New logo run-rate\n+ weighted pipeline\nACV by tier"),
    ("Q3: mart_pl (T4Q)", "OpEx ratios:\nCOGS%, S&M%,\nR&D%, G&A%"),
]
for i, (name, desc) in enumerate(queries):
    q_box = add_shape(slide, Inches(0.5 + i * 2.6), Inches(5.0), Inches(2.35), Inches(1.0), CARD_BG, ACCENT_BLUE)
    add_multi_text(q_box, [name, desc], font_size=9, color=LIGHT_GRAY, bold_first=True)

# Arrow to output
add_right_arrow(slide, Inches(5.3), Inches(6.35) - Inches(0.9), Inches(0.4), Inches(0.3), ACCENT_GREEN)
# Roll forward box
rf_box = add_shape(slide, Inches(5.9), Inches(5.0), Inches(2.0), Inches(1.0), ACCENT_GREEN)
set_shape_text(rf_box, "Roll Forward\n12 months\nx 3 tiers\n+ scenarios", font_size=10, color=WHITE, bold=True)


# ════════════════════════════════════════════════════════════════
# SLIDE 10 — SECURITY & IAM
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(8), Inches(0.6),
             "Security & IAM Model", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.5), Inches(0.9), Inches(1), Pt(3), ACCENT_RED)
add_text_box(slide, Inches(0.5), Inches(1.0), Inches(8), Inches(0.4),
             "LEVEL 300  |  Zero-password architecture with least-privilege IAM", font_size=12, color=ACCENT_RED, bold=True)

security_cards = [
    ("IAM Authentication", ACCENT_BLUE, [
        "No stored passwords anywhere",
        "Redshift: IAM-based auth via Data API",
        "Auto-provisioned DB users:",
        "  IAMR:acme-finance-dev-*-lambda",
        "S3: OAC + sigv4 (no public access)",
        "CloudFront: signed origins only",
        "Lambda: Function URL (AuthType=NONE",
        "  behind CloudFront — no API keys)"
    ]),
    ("Least Privilege", ACCENT_TEAL, [
        "Each Lambda has dedicated IAM role",
        "Roles scoped to specific actions:",
        "  text_to_sql: redshift-data:*",
        "  forecast: redshift-data:*",
        "  api: redshift-data + bedrock",
        "Redshift GRANT: SELECT only on marts",
        "  (no INSERT/UPDATE/DELETE)",
        "Glue: scoped to specific databases"
    ]),
    ("Data Protection", ACCENT_GREEN, [
        "S3: SSE-KMS encryption at rest",
        "Redshift: encrypted workgroup",
        "All data in single region (us-east-1)",
        "VPC: Redshift in private subnets",
        "  (NAT gateway for outbound)",
        "No PII in synthetic data",
        "Glue connection: encrypted JDBC",
        ""
    ]),
    ("Network Security", ACCENT_AMBER, [
        "CloudFront → Lambda (AWS backbone)",
        "Lambda → Redshift (VPC endpoint)",
        "Lambda → Bedrock (AWS API)",
        "Lambda → S3 (VPC endpoint)",
        "No public Redshift endpoint",
        "Security groups: Lambda → Redshift",
        "  only on port 5439",
        ""
    ]),
]

for i, (title, color, items) in enumerate(security_cards):
    x = Inches(0.5 + i * 3.15)
    hdr = add_shape(slide, x, Inches(1.7), Inches(2.9), Inches(0.45), color)
    set_shape_text(hdr, title, font_size=12, color=WHITE, bold=True)
    body = add_shape(slide, x, Inches(2.15), Inches(2.9), Inches(2.8), CARD_BG, color)
    add_multi_text(body, items, font_size=9, color=LIGHT_GRAY)

# IaC callout
iac = add_shape(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.0), CARD_BG, RGBColor(0xA7, 0x8B, 0xFA))
add_multi_text(iac, [
    "INFRASTRUCTURE AS CODE — Everything deployed via Terraform",
    "8 modules: networking, redshift, glue, bedrock, api, frontend, pipelines, monitoring",
    "State: S3 backend with DynamoDB locking  |  Environments: dev (prod follows same modules)",
    "CI/CD: GitHub Actions → terraform plan → manual approve → terraform apply"
], font_size=10, color=LIGHT_GRAY, bold_first=True)


# ════════════════════════════════════════════════════════════════
# SLIDE 11 — KEY DECISIONS & TRADE-OFFS
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(8), Inches(0.6),
             "Key Architectural Decisions", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.5), Inches(0.9), Inches(1), Pt(3), ACCENT_AMBER)

decisions = [
    ("Redshift Serverless\n(not Provisioned)", "Pay per query, auto-pause when idle. Lab/startup workload\ndoesn't justify 24/7 provisioned cluster. RPU auto-scales 8-32.", ACCENT_BLUE),
    ("Apache Iceberg\n(not raw Parquet)", "ACID writes, schema evolution, time travel. Spectrum reads\nIceberg natively — zero-copy from S3 to Redshift queries.", ACCENT_TEAL),
    ("Lambda Web Adapter\n(not API Gateway)", "Function URL + RESPONSE_STREAM = unbuffered SSE for chat.\nLWA translates invocations to HTTP — zero code changes local/cloud.", ACCENT_GREEN),
    ("dbt (not stored procs)", "Version-controlled SQL, built-in testing (41 tests),\nlineage graph, incremental materializations for marts.", ACCENT_AMBER),
    ("Glue ETL (not EMR)", "Serverless Spark — no cluster management. 2 workers for ~22MB\nraw data. Auto-scales if data grows. Integrated with Glue Catalog.", RGBColor(0xA7, 0x8B, 0xFA)),
    ("Driver-Based Forecast\n(not OLS trend)", "ARR bridge captures SaaS dynamics (churn/expansion/new logo)\nper tier. Scenario overrides enable what-if analysis.", ACCENT_RED),
    ("Bedrock Agent\n(not custom RAG)", "Managed agent orchestration with 6 action groups. No\nvector DB to maintain. Claude Sonnet for high-quality NL responses.", ACCENT_BLUE),
    ("Synthetic Data\n(not anonymized prod)", "Full control over data characteristics. Seeded anomalies for\ntesting. Realistic SaaS growth trajectory (FY23-25).", ACCENT_TEAL),
]

for i, (decision, rationale, color) in enumerate(decisions):
    row = i // 2
    col = i % 2
    x = Inches(0.5 + col * 6.3)
    y = Inches(1.4 + row * 1.45)

    d_box = add_shape(slide, x, y, Inches(6.0), Inches(1.25), CARD_BG, color)

    tf = d_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(6)

    p = tf.paragraphs[0]
    p.text = decision.replace("\n", " ")
    p.font.size = Pt(11)
    p.font.color.rgb = color
    p.font.bold = True

    for line in rationale.split("\n"):
        p2 = tf.add_paragraph()
        p2.text = line
        p2.font.size = Pt(9)
        p2.font.color.rgb = LIGHT_GRAY
        p2.space_before = Pt(2)


# ════════════════════════════════════════════════════════════════
# SLIDE 12 — SUMMARY / CLOSING
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_shape(slide, Inches(1), Inches(1.8), Inches(1.5), Pt(4), ACCENT_BLUE)

add_text_box(slide, Inches(1), Inches(2.1), Inches(11), Inches(0.9),
             "Architecture Summary", font_size=40, color=WHITE, bold=True)

summary_items = [
    ("3 Source Systems", "ERP + CRM + EPM with shared UUIDs for reconciliation"),
    ("Medallion Architecture", "S3 Raw (Bronze) → Iceberg Curated (Silver) → dbt Marts (Gold)"),
    ("Serverless-First", "Redshift Serverless, Glue ETL, Lambda, Step Functions — pay-per-use"),
    ("23 dbt Models, 41 Tests", "Full lineage, GL balance checks, P&L identity assertions"),
    ("AI-Powered Analytics", "Bedrock Agent with 6 action groups — NL to insights in seconds"),
    ("Driver-Based Forecasting", "ARR cohort model with tier-specific rates and scenario overrides"),
]

for i, (title, desc) in enumerate(summary_items):
    y = Inches(3.2 + i * 0.6)
    # Bullet
    bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), y + Inches(0.08), Inches(0.15), Inches(0.15))
    bullet.fill.solid()
    bullet.fill.fore_color.rgb = ACCENT_BLUE
    bullet.line.fill.background()
    bullet.shadow.inherit = False

    add_text_box(slide, Inches(1.6), y, Inches(3.5), Inches(0.35),
                 title, font_size=16, color=WHITE, bold=True)
    add_text_box(slide, Inches(5.2), y, Inches(7), Inches(0.35),
                 desc, font_size=14, color=LIGHT_GRAY)

add_text_box(slide, Inches(1), Inches(6.7), Inches(11), Inches(0.5),
             "Questions?", font_size=28, color=ACCENT_TEAL, bold=True,
             alignment=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════
# Save
# ════════════════════════════════════════════════════════════════
out_path = os.path.join(os.path.dirname(__file__), "ACME_Finance_Data_Journey.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
